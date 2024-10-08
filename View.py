# -*- coding: utf-8 -*-
"""
Created on Thu May 23 11:52:49 2024

@author: mmorin
"""

import tkinter as tk
from tkinter import ttk
from tkcalendar import DateEntry
from Data import DataManagement
from Model import PortfolioManagement
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.chart import XL_LABEL_POSITION
import win32com.client
import os
import pandas as pd
import matplotlib.pyplot as plt
from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg

class FinancialReportingApp(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("Financial Reporting")
        self.geometry("1000x800")
        self.configure(bg="#f5f5f5")
        self.data_manager = DataManagement()
        self.selected_companies = []
        self.start_date = tk.StringVar(value="01/10/2019")
        self.end_date = tk.StringVar(value="29/12/2019")
        self.create_widgets()
        
    def create_widgets(self):
        title_frame = tk.Frame(self, bg="#003366")
        title_frame.grid(row=0, column=0, columnspan=4, sticky="nsew", pady=10)
        tk.Label(title_frame, text="FINANCIAL REPORTING", font=("Arial", 14, "bold"), fg="white", bg="#003366").pack(pady=5)
        tk.Label(title_frame, text="Réalisé par Maxime Miriot-Jaubert, Maxence Morin, Mahak Shabeer", font=("Arial", 12), fg="white", bg="#003366").pack(pady=5)
        notebook = ttk.Notebook(self)
        notebook.grid(row=1, column=0, columnspan=4, sticky="nsew", pady=10, padx=10)
        self.create_determination_tab(notebook)
        self.create_overview_tab(notebook)
        self.create_visualization_tab(notebook)
        for row in range(3):
            self.grid_rowconfigure(row, weight=1, pad=10)
        for col in range(4):
            self.grid_columnconfigure(col, weight=1, pad=10)

    def create_determination_tab(self, notebook):
        determination_frame = ttk.Frame(notebook)
        notebook.add(determination_frame, text="Détermination du portefeuille")
        section_frame = tk.Frame(determination_frame, bg="#d9d9d9")
        section_frame.pack(fill="both", expand=True, pady=10)
        tk.Label(section_frame, text="Détermination du portefeuille", font=("Arial", 12, "bold"), bg="#d9d9d9").pack(pady=5)
        self.create_date_selectors(section_frame)
        company_portfolio_frame = tk.Frame(section_frame, bg="#d9d9d9")
        company_portfolio_frame.pack(fill="both", expand=True, pady=10)
        self.create_company_list(company_portfolio_frame)
        self.create_selected_portfolio(company_portfolio_frame)
        button_frame = tk.Frame(determination_frame, bg="#d9d9d9")
        button_frame.pack(pady=10, side=tk.BOTTOM)
        validate_button = tk.Button(button_frame, text="Valider portefeuille", command=self.update_list_isin, font=("Arial", 12, "bold"), bg="#007acc", fg="white", bd=0, width=18)
        validate_button.grid(row=0, column=0, padx=5)
        generate_data_button = tk.Button(button_frame, text="Générer les données du portefeuille", command=self.generate_portfolio_data, font=("Arial", 12, "bold"), bg="#007acc", fg="white", bd=0, width=28)
        generate_data_button.grid(row=0, column=1, padx=5)
        generate_ppt_button = tk.Button(button_frame, text="Générer Powerpoint", command=self.generate_powerpoint, font=("Arial", 12, "bold"), bg="#007acc", fg="white", bd=0, width=18)
        generate_ppt_button.grid(row=0, column=2, padx=5)
        generate_pdf_button = tk.Button(button_frame, text="Générer PDF", command=self.generate_pdf, font=("Arial", 12, "bold"), bg="#007acc", fg="white", bd=0, width=18)
        generate_pdf_button.grid(row=0, column=3, padx=5)

    def create_overview_tab(self, notebook):
        overview_frame = ttk.Frame(notebook)
        notebook.add(overview_frame, text="Overview et composition du portefeuille")
        # add a label
        tk.Label(overview_frame, text="Overview", font=("Arial", 14, "bold")).pack(pady=10)
        # old table
        self.overview_tree = ttk.Treeview(overview_frame)
        self.overview_tree.pack(expand=True, fill='both')
        # add a label
        tk.Label(overview_frame, text="Composition du portefeuille", font=("Arial", 14, "bold")).pack(pady=10)
        # add a table
        self.new_overview_tree = ttk.Treeview(overview_frame)
        self.new_overview_tree.pack(expand=True, fill='both')

    def display_portfolio_overview(self):
        # Obtenir liste_isin, start_date, end_date à partir de generate_list_isin
        liste_isin, start_date, end_date = self.generate_list_isin()
        # Instancier PortfolioManagement
        portfolio_management = PortfolioManagement()
        # Utiliser l'instance pour obtenir les données de rapport consolidées trimestrielles
        consolidated_quarter_reporting_data = portfolio_management.get_consolidated_quarter_reporting(liste_isin, start_date, end_date)
        # Supposons que vous voulez afficher les données dans un tableau ou un graphique
        # Création d'un tableau dans l'onglet d'aperçu
        columns = list(consolidated_quarter_reporting_data.columns)
        self.overview_tree["columns"] = columns
        self.overview_tree["show"] = "headings"  # Show only the headings
        # Définir les en-têtes de colonnes
        for col in columns:
            self.overview_tree.heading(col, text=col)
            self.overview_tree.column(col, width=100, anchor='center')
        # Supprimer les anciennes données du tableau
        for item in self.overview_tree.get_children():
            self.overview_tree.delete(item)
        # Ajouter les nouvelles données du portefeuille au tableau
        for index, row in consolidated_quarter_reporting_data.iterrows():
            self.overview_tree.insert("", "end", values=list(row))
        # Obtenir les données pour le nouveau tableau
        df_portfolio, optimal_weights, cumulative_returns = portfolio_management.reporting(liste_isin, start_date, end_date)
        # Préparer les données pour l'affichage
        df_features = self.data_manager.get_CAC40_stock_historical_price(features=True)
        isin_dict = df_features.set_index('ISIN_Portfolio')['Company_name'].to_dict()
        company_names = [isin_dict[isin] for isin in liste_isin]
        reporting_data = pd.DataFrame({
            'ISIN': liste_isin,
            'Company': company_names,
            'Weight': optimal_weights.values,
            'Cumulative Return': cumulative_returns.iloc[-1].values
        })
        # Création du nouveau tableau dans l'onglet d'aperçu
        new_columns = list(reporting_data.columns)
        self.new_overview_tree["columns"] = new_columns
        self.new_overview_tree["show"] = "headings"  # Show only the headings
        # Définir les en-têtes de colonnes pour le nouveau tableau
        for col in new_columns:
            self.new_overview_tree.heading(col, text=col)
            self.new_overview_tree.column(col, width=100, anchor='center')
        # Supprimer les anciennes données du nouveau tableau
        for item in self.new_overview_tree.get_children():
            self.new_overview_tree.delete(item)
        # Ajouter les nouvelles données du portefeuille au nouveau tableau
        for index, row in reporting_data.iterrows():
            self.new_overview_tree.insert("", "end", values=list(row))

    def create_visualization_tab(self, notebook):
        visualization_frame = ttk.Frame(notebook)
        notebook.add(visualization_frame, text="Visualisation graphique du portefeuille et de son benchmark")
        tk.Label(visualization_frame, text="Graphiques et visualisations", font=("Arial", 14, "bold")).pack(pady=10)
        # Ajouter un cadre pour les graphiques
        self.chart_frame = tk.Frame(visualization_frame, bg="#f5f5f5")
        self.chart_frame.pack(fill="both", expand=True, pady=10)
        # Ajouter le premier graphe
        self.chart1 = tk.Canvas(self.chart_frame, bg="white")
        self.chart1.pack(fill="both", expand=True, padx=10, pady=(10, 5))  # Ajouter un espacement autour du graphe
        # Ajouter un espace entre les graphes
        spacer = tk.Frame(self.chart_frame, height=76, bg="#f5f5f5")  # Ajouter un espacement vertical de 2 cm
        spacer.pack(fill="x", pady=(50, 50))
        # Ajouter le second graphe
        self.chart2 = tk.Canvas(self.chart_frame, bg="white")
        self.chart2.pack(fill="both", expand=True, padx=10, pady=(5, 10))  # Ajouter un espacement autour du graphe


    def plot_portfolio_graph(self, ax):
        portfolio_management = PortfolioManagement()
        df_portfolio, _, _ = portfolio_management.reporting(self.generate_list_isin()[0], self.start_date.get(), self.end_date.get())
        if isinstance(df_portfolio, tuple):
            df_portfolio = df_portfolio[0]  # Access the DataFrame if it's a part of a tuple
        ax.plot(df_portfolio.index, df_portfolio["VL"], label="Portefeuille")
        ax.set_ylabel("Valeur Portefeuille choisi")
        ax.legend()

    def plot_benchmark_graph(self, ax):
        # Obtenir les dates de début et de fin
        start_date = self.start_date.get()
        end_date = self.end_date.get()
        # Utiliser la bonne méthode pour obtenir les données historiques du CAC40
        df_historical_price_CAC40 = self.data_manager.get_CAC40_index_historical_price(start_date, end_date)
        # Assurez-vous que la colonne 'Adj Close' est présente
        if 'Adj Close' not in df_historical_price_CAC40.columns:
            raise KeyError("La colonne 'Adj Close' n'existe pas dans le DataFrame.")
        ax.plot(df_historical_price_CAC40.index, df_historical_price_CAC40["Adj Close"], label="CAC40")
        ax.set_xlabel("Temps")
        ax.set_ylabel("Valeur CAC 40")
        ax.legend()

    def create_date_selectors(self, parent):
        frame = tk.Frame(parent, bg="#d9d9d9")
        frame.pack(pady=5)
        tk.Label(frame, text="Date de début:", font=("Arial", 10), bg="#d9d9d9").grid(row=0, column=0, padx=5, pady=5, sticky="e")
        start_date_entry = DateEntry(frame, textvariable=self.start_date, date_pattern='dd/mm/yyyy', font=("Arial", 10), year=2019, month=10, day=1)
        start_date_entry.grid(row=0, column=1, padx=5, pady=5)
        tk.Label(frame, text="Date de fin:", font=("Arial", 10), bg="#d9d9d9").grid(row=0, column=2, padx=5, pady=5, sticky="e")
        end_date_entry = DateEntry(frame, textvariable=self.end_date, date_pattern='dd/mm/yyyy', font=("Arial", 10), year=2019, month=12, day=29)
        end_date_entry.grid(row=0, column=3, padx=5, pady=5)

    def create_company_list(self, parent):
        frame = tk.Frame(parent, bg="#ffffff", bd=2, relief=tk.GROOVE)
        frame.pack(side="left", fill="both", expand=True, padx=10, pady=10)
        tk.Label(frame, text="Liste des entreprises", font=("Arial", 12, "bold"), bg="#ffffff").pack(pady=5)
        df_features = self.data_manager.get_CAC40_stock_historical_price(features=True)
        companies = df_features['Company_name'].dropna().tolist()
        self.company_listbox = tk.Listbox(frame, selectmode=tk.MULTIPLE, height=10, font=("Arial", 10))
        for company in companies:
            self.company_listbox.insert(tk.END, company)
        self.company_listbox.pack(fill="both", expand=True, padx=5, pady=5)
        add_button = tk.Button(frame, text="Ajouter", command=self.add_selected_companies, font=("Arial", 10), bg="#007acc", fg="white", bd=0)
        add_button.pack(pady=5)

    def create_selected_portfolio(self, parent):
        frame = tk.Frame(parent, bg="#ffffff", bd=2, relief=tk.GROOVE)
        frame.pack(side="right", fill="both", expand=True, padx=10, pady=10)

        tk.Label(frame, text="Portefeuille sélectionné", font=("Arial", 12, "bold"), bg="#ffffff").pack(pady=5)

        scrollbar = tk.Scrollbar(frame)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)

        self.selected_companies_listbox = tk.Listbox(frame, yscrollcommand=scrollbar.set, height=10, font=("Arial", 10))
        self.selected_companies_listbox.pack(fill="both", expand=True, padx=5, pady=5)
        scrollbar.config(command=self.selected_companies_listbox.yview)

        remove_button = tk.Button(frame, text="Supprimer", command=self.remove_selected_companies, font=("Arial", 10), bg="#ff1a1a", fg="white", bd=0)
        remove_button.pack(pady=5)

    def add_selected_companies(self):
        selected_indices = self.company_listbox.curselection()
        selected_companies = [self.company_listbox.get(i) for i in selected_indices]
        for company in selected_companies:
            if company not in self.selected_companies:
                self.selected_companies.append(company)
        self.update_selected_companies()

    def remove_selected_companies(self):
        selected_indices = self.selected_companies_listbox.curselection()
        selected_companies = [self.selected_companies_listbox.get(i) for i in selected_indices]
        for company in selected_companies:
            self.selected_companies.remove(company)
        self.update_selected_companies()

    def update_selected_companies(self):
        self.selected_companies_listbox.delete(0, tk.END)
        for company in self.selected_companies:
            self.selected_companies_listbox.insert(tk.END, company)

    def update_list_isin(self):
        list_isin, start_date, end_date = self.generate_list_isin()
        print(f"List ISIN updated: {list_isin}")
        print(f"Start date: {start_date}")
        print(f"End date: {end_date}")

    def generate_list_isin(self):
        df_features = self.data_manager.get_CAC40_stock_historical_price(features=True)
        isin_dict = df_features.set_index('Company_name')['ISIN_Portfolio'].to_dict()
        list_isin = [isin_dict[company] for company in self.selected_companies]
        return list_isin, self.start_date.get(), self.end_date.get()

    def generate_portfolio_data(self):
        self.display_portfolio_overview()
        self.display_portfolio_visualization()

    def display_portfolio_visualization(self):
        fig, (ax1, ax2) = plt.subplots(2, 1, figsize=(10, 8))
        self.plot_portfolio_graph(ax1)
        self.plot_benchmark_graph(ax2)
        for widget in self.chart_frame.winfo_children():
            widget.destroy()
        canvas = FigureCanvasTkAgg(fig, master=self.chart_frame)
        canvas.draw()
        canvas.get_tk_widget().pack(side=tk.TOP, fill=tk.BOTH, expand=True)


    def save_charts_as_images(self):
        # Create a single figure with two subplots
        fig, (ax1, ax2) = plt.subplots(2, 1, figsize=(10, 8))
        # Plot the graphs in the same figure
        self.plot_portfolio_graph(ax1)
        self.plot_benchmark_graph(ax2)
        # Save the combined figure as a single image
        combined_image_path = "combined_graphs.png"
        fig.savefig(combined_image_path)
        plt.close(fig)
        return combined_image_path

    def generate_powerpoint(self):
        # Créer une nouvelle présentation PowerPoint
        prs = Presentation()
        # Ajouter une diapositive de titre
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Financial Reporting"
        subtitle.text = "Overview et composition du portefeuille\nVisualisation graphique du portefeuille et de son benchmark"
        # Ajouter une diapositive pour l'overview et la composition du portefeuille
        overview_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(overview_slide_layout)
        title = slide.shapes.title
        title.text = "Overview et composition du portefeuille"
        # Obtenir les dimensions de la diapositive
        slide_width = prs.slide_width
        slide_height = prs.slide_height
    
        # Ajouter les éléments des Treeviews à la diapositive "Overview et composition du portefeuille"
        def add_treeview_to_slide(treeview, slide, start_top):
            rows = [list(treeview.heading(col)["text"] for col in treeview["columns"])]
            rows += [list(treeview.item(item)["values"]) for item in treeview.get_children()]
            # Calculer la hauteur nécessaire pour le tableau
            row_height = Pt(18)  # Hauteur estimée de chaque ligne
            table_height = row_height * len(rows)
            table_width = slide_width - Inches(1)  # Largeur totale moins marges
            # Créer un tableau dans la diapositive
            table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(0.5), start_top, table_width, table_height).table
            # Remplir le tableau avec les données
            for i, row in enumerate(rows):
                for j, val in enumerate(row):
                    table.cell(i, j).text = str(val)
            # Ajuster la taille de la police pour adapter le texte dans les cellules
            for i in range(len(rows)):
                for j in range(len(rows[0])):
                    cell = table.cell(i, j)
                    cell.text_frame.text = str(rows[i][j])
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(10)
        # Ajouter le premier Treeview
        add_treeview_to_slide(self.overview_tree, slide, Inches(1.5))
        # Ajouter le second Treeview
        add_treeview_to_slide(self.new_overview_tree, slide, Inches(3.5))
        # Ajouter une diapositive pour la visualisation graphique
        visualization_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(visualization_slide_layout)
        title = slide.shapes.title
        title.text = "Visualisation graphique du portefeuille et de son benchmark"
        # Enregistrer les graphiques en tant qu'image
        combined_image_path = self.save_charts_as_images()
        # Ajouter l'image à la diapositive
        slide.shapes.add_picture(combined_image_path, Inches(1), Inches(1.5), width=Inches(8), height=Inches(7))
        # Enregistrer la présentation
        prs.save('Financial_Reporting_Presentation.pptx')
        print("PowerPoint generated successfully.")
        # Créer une nouvelle présentation PowerPoint
        prs = Presentation()
        # Ajouter une diapositive de titre
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Financial Reporting"
        subtitle.text = "Overview et composition du portefeuille\nVisualisation graphique du portefeuille et de son benchmark"
        # Ajouter une diapositive pour l'overview et la composition du portefeuille
        overview_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(overview_slide_layout)
        title = slide.shapes.title
        title.text = "Overview et composition du portefeuille"
        # Obtenir les dimensions de la diapositive
        slide_width = prs.slide_width
        slide_height = prs.slide_height
    
        # Ajouter les éléments des Treeviews à la diapositive "Overview et composition du portefeuille"
        def add_treeview_to_slide(treeview, slide, start_top):
            rows = [list(treeview.heading(col)["text"] for col in treeview["columns"])]
            rows += [list(treeview.item(item)["values"]) for item in treeview.get_children()]
            # Calculer la hauteur nécessaire pour le tableau
            row_height = Pt(18)  # Hauteur estimée de chaque ligne
            table_height = row_height * len(rows)
            table_width = slide_width - Inches(1)  # Largeur totale moins marges
            # Créer un tableau dans la diapositive
            table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(0.5), start_top, table_width, table_height).table
            # Remplir le tableau avec les données
            for i, row in enumerate(rows):
                for j, val in enumerate(row):
                    table.cell(i, j).text = str(val) 
            # Ajuster la taille de la police pour adapter le texte dans les cellules
            for i in range(len(rows)):
                for j in range(len(rows[0])):
                    cell = table.cell(i, j)
                    cell.text_frame.text = str(rows[i][j])
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(10)
        # Ajouter le premier Treeview
        add_treeview_to_slide(self.overview_tree, slide, Inches(1.5))
        # Ajouter le second Treeview
        add_treeview_to_slide(self.new_overview_tree, slide, Inches(3.5))
        # Ajouter une diapositive pour la visualisation graphique
        visualization_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(visualization_slide_layout)
        title = slide.shapes.title
        title.text = "Visualisation graphique du portefeuille et de son benchmark"
        # Enregistrer les graphiques en tant qu'images
        portfolio_image_path, benchmark_image_path = self.save_charts_as_images()
        # Ajouter les images à la diapositive
        slide.shapes.add_picture(portfolio_image_path, Inches(1), Inches(1.5), width=Inches(8), height=Inches(3.5))
        slide.shapes.add_picture(benchmark_image_path, Inches(1), Inches(5), width=Inches(8), height=Inches(3.5))
        # Enregistrer la présentation
        prs.save('Financial_Reporting_Presentation.pptx')
        print("PowerPoint generated successfully.")

    def generate_pdf(self):
        # Chemin vers le fichier PowerPoint généré
        ppt_filename = 'Financial_Reporting_Presentation.pptx'
        pdf_filename = 'Financial_Reporting_Presentation.pdf'
        # Obtenir le chemin absolu des fichiers
        ppt_path = os.path.abspath(ppt_filename)
        pdf_path = os.path.abspath(pdf_filename)
        # Vérifiez si le fichier PowerPoint existe
        if not os.path.exists(ppt_path):
            print(f"Le fichier PowerPoint {ppt_path} n'existe pas.")
            return
        try:
            # Initialiser PowerPoint
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            powerpoint.Visible = True
            # Ouvrir la présentation PowerPoint
            presentation = powerpoint.Presentations.Open(ppt_path)
            # Enregistrer la présentation en tant que PDF
            presentation.SaveAs(pdf_path, 32)  # 32 correspond au format PDF
            # Fermer la présentation et PowerPoint
            presentation.Close()
            powerpoint.Quit()
            print("PDF generated successfully.")
        except Exception as e:
            print(f"An error occurred: {e}")